from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x0D, 0x1A)   # deep navy
ACCENT1 = RGBColor(0x00, 0xC8, 0xFF)   # electric cyan
ACCENT2 = RGBColor(0x7C, 0x3A, 0xFF)   # violet
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xB0, 0xB8, 0xC8)
CARD    = RGBColor(0x1A, 0x1E, 0x35)   # card surface

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ── Helpers ───────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_bullet_card(slide, x, y, w, h, title, bullets):
    """Dark card with title + bullet points."""
    card = add_rect(slide, x, y, w, h, CARD)
    # top accent bar
    add_rect(slide, x, y, w, Inches(0.07), ACCENT1)
    add_text(slide, title, x + Inches(0.2), y + Inches(0.12),
             w - Inches(0.4), Inches(0.45), size=16, bold=True, color=ACCENT1)
    bullet_y = y + Inches(0.65)
    for b in bullets:
        tb = slide.shapes.add_textbox(x + Inches(0.2), bullet_y,
                                      w - Inches(0.4), Inches(0.45))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "▸  " + b
        run.font.size = Pt(13)
        run.font.color.rgb = LGRAY
        bullet_y += Inches(0.42)

def add_stat_card(slide, x, y, w, h, stat, label):
    add_rect(slide, x, y, w, h, CARD)
    add_rect(slide, x, y, w, Inches(0.07), ACCENT2)
    add_text(slide, stat, x, y + Inches(0.2), w, Inches(0.8),
             size=40, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.95), w, Inches(0.4),
             size=12, color=LGRAY, align=PP_ALIGN.CENTER)

def title_pill(slide, text, x, y):
    """Small glowing pill label."""
    add_rect(slide, x, y, Inches(2.8), Inches(0.32), ACCENT2)
    add_text(slide, text, x, y + Inches(0.02), Inches(2.8), Inches(0.32),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
# Left gradient accent strip
add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, ACCENT1)
add_rect(slide, Inches(0.18), Inches(0), Inches(0.06), H, ACCENT2)
# Main heading
add_text(slide, "HYBRID", Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "WEB APPLICATION FIREWALL",
         Inches(0.5), Inches(2.4), Inches(12), Inches(0.9),
         size=36, bold=True, color=ACCENT1, align=PP_ALIGN.LEFT)
add_rect(slide, Inches(0.5), Inches(3.4), Inches(6), Inches(0.05), ACCENT2)
add_text(slide, "Intelligent Threat Detection via Signature Matching & Machine Learning",
         Inches(0.5), Inches(3.55), Inches(10), Inches(0.5),
         size=16, color=LGRAY, italic=True)
add_text(slide, "Submitted by: [Team Names]  |  Guided by: Dr. [Guide Name]",
         Inches(0.5), Inches(6.6), Inches(10), Inches(0.4),
         size=13, color=LGRAY)
add_text(slide, "Dept. of CSE  |  Academic Year 2025–26",
         Inches(0.5), Inches(7.0), Inches(10), Inches(0.35),
         size=12, color=ACCENT2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), CARD)
add_text(slide, "AGENDA", Inches(0.5), Inches(0.2), W, Inches(0.7),
         size=34, bold=True, color=WHITE)
add_rect(slide, Inches(0.5), Inches(1.0), Inches(0.06), Inches(6), ACCENT1)

items = [
    "01  Problem Statement & Motivation",
    "02  Proposed Hybrid-WAF Architecture",
    "03  Signature Engine & Rule System",
    "04  Machine Learning Module (Isolation Forest)",
    "05  Fusion Layer & Decision Logic",
    "06  Implementation Stack & Live Dashboard",
    "07  Results & Performance Metrics",
    "08  Conclusion & Future Scope",
]
for i, item in enumerate(items):
    c = ACCENT1 if i % 2 == 0 else ACCENT2
    add_text(slide, item, Inches(0.8), Inches(1.15) + Inches(0.65) * i,
             Inches(11), Inches(0.55), size=17, color=WHITE if i % 2 == 0 else LGRAY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), CARD)
title_pill(slide, "PROBLEM STATEMENT", Inches(0.5), Inches(0.2))
add_text(slide, "The Detection Gap", Inches(0.5), Inches(0.55), Inches(10), Inches(0.55),
         size=30, bold=True, color=WHITE)

cards = [
    ("Signature-Based WAF", ["✓ Fast & precise for known threats",
                              "✗ Blind to zero-day exploits",
                              "✗ Defeated by polymorphic payloads"]),
    ("Anomaly-Based Systems", ["✓ Detects novel/unknown attacks",
                               "✗ High false-positive rates",
                               "✗ Disrupts legitimate traffic"]),
    ("The Gap", ["No single approach is sufficient",
                 "Attackers exploit the blind spots",
                 "Need: Unified intelligent WAF"]),
]
cx = Inches(0.4)
for (title, bullets) in cards:
    add_bullet_card(slide, cx, Inches(1.35), Inches(4.1), Inches(5.8), title, bullets)
    cx += Inches(4.3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), CARD)
title_pill(slide, "SYSTEM DESIGN", Inches(0.5), Inches(0.2))
add_text(slide, "Hybrid-WAF Architecture", Inches(0.5), Inches(0.55), Inches(10), Inches(0.55),
         size=30, bold=True, color=WHITE)

flow = [
    ("Network Traffic", ACCENT2),
    ("Capture Engine\n(Scapy)", CARD),
    ("Feature Extractor\n(19 Features)", CARD),
    ("Signature Engine", ACCENT1),
    ("ML Module", ACCENT2),
    ("Fusion Layer\n(Risk Score)", rgb(0, 180, 80)),
    ("Block / Allow", CARD),
]
bx = Inches(0.25)
for i, (label, col) in enumerate(flow):
    add_rect(slide, bx, Inches(2.8), Inches(1.75), Inches(1.2), col)
    add_text(slide, label, bx, Inches(2.85), Inches(1.75), Inches(1.1),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        add_text(slide, "▶", bx + Inches(1.75), Inches(3.1), Inches(0.35), Inches(0.6),
                 size=18, color=ACCENT1, align=PP_ALIGN.CENTER)
    bx += Inches(1.9)

add_text(slide, "Parallel dual-engine analysis with intelligent Fusion Layer arbitration.",
         Inches(0.5), Inches(4.35), Inches(12), Inches(0.4),
         size=14, italic=True, color=LGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Signature Engine
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), CARD)
title_pill(slide, "ENGINE 1", Inches(0.5), Inches(0.2))
add_text(slide, "Signature Engine — Rule-Based Detection", Inches(0.5), Inches(0.55),
         Inches(12), Inches(0.55), size=26, bold=True, color=WHITE)

add_bullet_card(slide, Inches(0.4), Inches(1.35), Inches(6.1), Inches(5.8),
    "How it Works",
    ["Regex patterns loaded from YAML rule files.",
     "Rules have weighted severity scores (1–10).",
     "Cumulative score triggers Block decision.",
     "Supports multi-pattern collaborative matching.",
     "Targets: SQLi, XSS, Path Traversal, RFI."])

add_bullet_card(slide, Inches(6.8), Inches(1.35), Inches(6.1), Inches(2.7),
    "Attack Example — SQL Injection",
    ["Input: ' OR 1=1; DROP TABLE users--",
     "Rule Hit: UNION + DROP keywords",
     "Score: 9/10 → Instant Block"])
add_bullet_card(slide, Inches(6.8), Inches(4.25), Inches(6.1), Inches(2.9),
    "Attack Example — XSS",
    ["Input: <script>alert('xss')</script>",
     "Rule Hit: <script> tag pattern matched",
     "Score: 8/10 → Blocked"])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ML Module
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), CARD)
title_pill(slide, "ENGINE 2", Inches(0.5), Inches(0.2))
add_text(slide, "ML Module — Isolation Forest Anomaly Detector",
         Inches(0.5), Inches(0.55), Inches(12), Inches(0.55),
         size=24, bold=True, color=WHITE)

add_bullet_card(slide, Inches(0.4), Inches(1.35), Inches(6.1), Inches(5.8),
    "Why Isolation Forest?",
    ["Unsupervised: No labeled attack data needed.",
     "Isolates anomalies via random binary splitting.",
     "Short path length → High anomaly score.",
     "Linear time O(n) — ideal for real-time.",
     "Trained on CICIDS2017 (80 → 19 features)."])

add_bullet_card(slide, Inches(6.8), Inches(1.35), Inches(6.1), Inches(2.6),
    "Key Features Extracted",
    ["Flow IAT Mean & Std (bot detection)",
     "Bwd Packet Length (exfiltration)",
     "SYN/ACK/PSH Flags (scanning)"])
add_bullet_card(slide, Inches(6.8), Inches(4.15), Inches(6.1), Inches(3.0),
    "Model Configuration",
    ["Estimators: 100 Trees",
     "Contamination Factor: 0.05",
     "Sliding Window: Last 1000 flows",
     "Score: 0.0 (safe) → 1.0 (anomalous)"])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Fusion Layer
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), CARD)
title_pill(slide, "DECISION LOGIC", Inches(0.5), Inches(0.2))
add_text(slide, "The Fusion Layer — Intelligent Arbitration",
         Inches(0.5), Inches(0.55), Inches(12), Inches(0.55),
         size=26, bold=True, color=WHITE)

add_rect(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(1.0), CARD)
add_rect(slide, Inches(1.5), Inches(1.6), Inches(0.08), Inches(1.0), ACCENT1)
add_text(slide, "Risk Score  =  (0.55 × Rule Score)  +  (0.45 × ML Score)",
         Inches(1.7), Inches(1.7), Inches(10), Inches(0.75),
         size=22, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

logic = [
    ("Signature CRITICAL", "Immediate Block — no ML check needed", ACCENT1),
    ("Score > 0.75", "Block — High confidence from both engines", rgb(220, 60, 60)),
    ("Score 0.45 – 0.75", "Flag & Alert — Human review recommended", rgb(255, 160, 0)),
    ("Score < 0.45", "Allow — Benign traffic pattern detected", rgb(0, 200, 100)),
]
row_y = Inches(2.85)
for label, desc, col in logic:
    add_rect(slide, Inches(0.5), row_y, Inches(3.0), Inches(0.65), col)
    add_text(slide, label, Inches(0.5), row_y + Inches(0.1), Inches(3.0), Inches(0.5),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, Inches(3.7), row_y + Inches(0.12), Inches(8.5), Inches(0.45),
             size=14, color=LGRAY)
    row_y += Inches(0.88)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Tech Stack & Dashboard
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), CARD)
title_pill(slide, "IMPLEMENTATION", Inches(0.5), Inches(0.2))
add_text(slide, "Technology Stack & Live Dashboard",
         Inches(0.5), Inches(0.55), Inches(12), Inches(0.55),
         size=26, bold=True, color=WHITE)

stack = [
    ("🐍  Python 3.10", "Core logic & orchestration"),
    ("📡  Scapy", "Live packet sniffing"),
    ("🤖  Scikit-Learn", "Isolation Forest ML model"),
    ("⚡  Pandas / NumPy", "High-speed feature extraction"),
    ("📊  Streamlit", "Real-time security dashboard"),
    ("📁  PyYAML", "Rule configuration management"),
]
sx = Inches(0.4)
sy = Inches(1.45)
for i, (tool, desc) in enumerate(stack):
    col = Inches(0.4) + (i % 3) * Inches(4.35)
    row = Inches(1.45) + (i // 3) * Inches(1.5)
    add_rect(slide, col, row, Inches(4.0), Inches(1.2), CARD)
    add_rect(slide, col, row, Inches(4.0), Inches(0.06), ACCENT2)
    add_text(slide, tool, col + Inches(0.15), row + Inches(0.1), Inches(3.7), Inches(0.5),
             size=15, bold=True, color=WHITE)
    add_text(slide, desc, col + Inches(0.15), row + Inches(0.6), Inches(3.7), Inches(0.45),
             size=13, color=LGRAY)

add_bullet_card(slide, Inches(0.4), Inches(4.6), Inches(12.4), Inches(2.6),
    "Dashboard Features",
    ["Live Traffic Ticker — real-time allow/block feed with IP, method, path.",
     "Threat Heatmap — visualizes attack density by source IP and time.",
     "AI Explainability — human-readable reasoning for every ML decision."])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Results
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), CARD)
title_pill(slide, "EVALUATION", Inches(0.5), Inches(0.2))
add_text(slide, "Results & Performance Metrics",
         Inches(0.5), Inches(0.55), Inches(12), Inches(0.55),
         size=26, bold=True, color=WHITE)

stats = [
    ("99.1%", "Detection Recall"),
    ("97.25%", "Precision"),
    ("0.28%", "False Positive Rate"),
    ("12ms", "Avg. Latency / Packet"),
]
sx = Inches(0.4)
for stat, label in stats:
    add_stat_card(slide, sx, Inches(1.4), Inches(3.0), Inches(1.6), stat, label)
    sx += Inches(3.2)

add_bullet_card(slide, Inches(0.4), Inches(3.2), Inches(6.1), Inches(4.0),
    "Security Benchmarking",
    ["Tested against: Slowloris DoS, SQLmap, Nmap scans.",
     "Detected 3/5 novel Zero-Day SQLi bypass attempts.",
     "Outperforms standalone ModSecurity for novel threats.",
     "Confusion Matrix: TP=4955, FP=140, TN=49860, FN=45."])
add_bullet_card(slide, Inches(6.8), Inches(3.2), Inches(6.1), Inches(4.0),
    "System Overhead",
    ["Memory: ~450 MB (model + flow buffers).",
     "CPU: ~15% on i7 at 100 req/sec load.",
     "12ms overhead vs 8ms (signature-only).",
     "Trade-off justified by detection gains."])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Future Scope
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), CARD)
title_pill(slide, "FUTURE SCOPE", Inches(0.5), Inches(0.2))
add_text(slide, "Limitations & Roadmap",
         Inches(0.5), Inches(0.55), Inches(12), Inches(0.55),
         size=28, bold=True, color=WHITE)

add_bullet_card(slide, Inches(0.4), Inches(1.35), Inches(6.1), Inches(5.8),
    "Current Limitations",
    ["No TLS decryption for HTTPS inspection.",
     "Training data bias from CICIDS2017.",
     "High memory for large flow windows.",
     "Single-node — no distributed support yet."])

roadmap = [
    ("Phase 1", "TLS inspection with root CA integration."),
    ("Phase 2", "Federated learning across WAF nodes."),
    ("Phase 3", "LLM-powered incident post-mortems."),
    ("Phase 4", "FPGA/SmartNIC hardware acceleration."),
]
ry = Inches(1.35)
for phase, desc in roadmap:
    add_rect(slide, Inches(6.8), ry, Inches(1.1), Inches(1.1), ACCENT2)
    add_text(slide, phase, Inches(6.8), ry + Inches(0.3), Inches(1.1), Inches(0.5),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(8.0), ry + Inches(0.25), Inches(4.9), Inches(0.6), CARD)
    add_text(slide, desc, Inches(8.1), ry + Inches(0.3), Inches(4.7), Inches(0.5),
             size=14, color=LGRAY)
    ry += Inches(1.4)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Thank You
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, ACCENT1)
add_rect(slide, Inches(0.18), Inches(0), Inches(0.06), H, ACCENT2)
add_text(slide, "THANK YOU", Inches(0.5), Inches(2.0), Inches(12), Inches(2.0),
         size=90, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(3), Inches(4.0), Inches(7.3), Inches(0.06), ACCENT1)
add_text(slide, "Questions & Discussion Welcome",
         Inches(0.5), Inches(4.2), Inches(12), Inches(0.6),
         size=22, italic=True, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Hybrid-WAF  ·  Intelligent Web Security  ·  CSE Dept. 2025–26",
         Inches(0.5), Inches(6.8), Inches(12), Inches(0.45),
         size=13, color=ACCENT2, align=PP_ALIGN.CENTER)


prs.save('Hybrid_WAF_Presentation_v2.pptx')
print("Done: Hybrid_WAF_Presentation_v2.pptx")
